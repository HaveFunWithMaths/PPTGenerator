"""
Chemistry PYQ PowerPoint Generator
Generates a presentation from chemistry previous year questions.

Author: Chemistry Teacher Assistant
Date: 2025
"""

from dataclasses import dataclass
from typing import List, Dict, Optional
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# CONFIGURATION & CONSTANTS
# =============================================================================

@dataclass(frozen=True)
class SlideColors:
    """Color scheme for presentation slides."""
    background: RGBColor = RGBColor(0, 0, 0)
    text: RGBColor = RGBColor(255, 255, 255)
    year: RGBColor = RGBColor(0, 212, 255)      # Vibrant Cyan
    marks: RGBColor = RGBColor(255, 179, 71)    # Golden Orange
    accent_line: RGBColor = RGBColor(60, 60, 80)
    question_badge: RGBColor = RGBColor(100, 100, 100)


@dataclass(frozen=True)
class SlideLayout:
    """Layout dimensions and margins for slides."""
    margin: float = 0.5          # inches
    top_margin: float = 0.3      # inches
    header_height: float = 0.4   # inches
    accent_line_y: float = 0.85  # inches
    accent_line_height: float = 0.03  # inches
    question_top: float = 1.1    # inches
    badge_bottom: float = 0.6    # inches
    content_ratio: float = 0.9   # 90% of slide height for content


@dataclass(frozen=True)
class FontSettings:
    """Typography settings."""
    name: str = "Calibri"
    question_size: int = 18
    meta_size: int = 16
    badge_size: int = 12


# Default instances
COLORS = SlideColors()
LAYOUT = SlideLayout()
FONTS = FontSettings()


# =============================================================================
# TEXT PROCESSING UTILITIES
# =============================================================================

# Unicode subscript and superscript mappings
SUBSCRIPT_MAP = str.maketrans("0123456789()+-=xyzn", "₀₁₂₃₄₅₆₇₈₉₍₎₊₋₌ₓyzₙ")
SUPERSCRIPT_MAP = str.maketrans("0123456789()+-=xyzn", "⁰¹²³⁴⁵⁶⁷⁸⁹⁽⁾⁺⁻⁼ˣʸᶻⁿ")

# LaTeX replacement patterns
LATEX_REPLACEMENTS = {
    # Arrows
    r"\rightarrow": "→",
    r"\longrightarrow": "→",
    r"\leftarrow": "←",
    r"\leftrightarrow": "↔",
    # Greek letters
    r"\Delta": "Δ",
    r"\underline{\Delta}": "Δ",
    r"\Omega": "Ω",
    r"\omega": "ω",
    r"\rho": "ρ",
    r"\pi": "π",
    r"\alpha": "α",
    r"\beta": "β",
    r"\gamma": "γ",
    r"\lambda": "λ",
    r"\mu": "μ",
    r"\sigma": "σ",
    r"\theta": "θ",
    r"\phi": "φ",
    r"\epsilon": "ε",
    r"\eta": "η",
    r"\tau": "τ",
    # Math operators
    r"\times": "×",
    r"\div": "÷",
    r"\pm": "±",
    r"\mp": "∓",
    r"\cdot": "·",
    r"\geq": "≥",
    r"\leq": "≤",
    r"\neq": "≠",
    r"\approx": "≈",
    r"\equiv": "≡",
    r"\propto": "∝",
    r"\infty": "∞",
    r"\sqrt": "√",
    # Other symbols
    r"\degree": "°",
    r"\circ": "°",
}


def _replace_subscript(match: re.Match) -> str:
    """Convert matched subscript pattern to Unicode subscript."""
    return match.group(1).translate(SUBSCRIPT_MAP)


def _replace_superscript(match: re.Match) -> str:
    """Convert matched superscript pattern to Unicode superscript."""
    return match.group(1).translate(SUPERSCRIPT_MAP)


def _replace_fraction(match: re.Match) -> str:
    """Convert matched fraction pattern to readable format."""
    numerator = match.group(1)
    denominator = match.group(2)
    return f"({numerator}/{denominator})"


def clean_chemistry_text(text: Optional[str]) -> str:
    """
    Convert LaTeX-style formatting to readable Unicode.
    
    Handles:
        - Subscripts (_{...} and _x patterns)
        - Superscripts (^{...} and ^x patterns)
        - Fractions (\frac{a}{b})  
        - Chemical arrows
        - Greek letters (Delta, Omega, etc.)
        - Math operators
        - Dollar sign delimiters
    
    Args:
        text: Raw text potentially containing LaTeX formatting
        
    Returns:
        Cleaned text with Unicode symbols
    """
    if not text:
        return ""
    
    # Convert fractions first: \frac{a}{b} -> (a/b)
    text = re.sub(r"\\frac\{([^}]+)\}\{([^}]+)\}", _replace_fraction, text)
    
    # Apply LaTeX symbol replacements
    for latex, unicode_char in LATEX_REPLACEMENTS.items():
        text = text.replace(latex, unicode_char)
    
    # Remove $ delimiters
    text = text.replace("$", "")
    
    # Convert superscripts: ^{...} pattern
    text = re.sub(r"\^\{([^}]+)\}", _replace_superscript, text)
    
    # Convert single character superscripts: ^2 pattern
    text = re.sub(r"\^([0-9])", _replace_superscript, text)
    
    # Convert subscripts: _{...} pattern
    text = re.sub(r"_\{([^}]+)\}", _replace_subscript, text)
    
    # Convert single character subscripts: _2 pattern
    text = re.sub(r"_([0-9])", _replace_subscript, text)
    
    # Clean up remaining backslashes from unknown LaTeX commands
    text = re.sub(r"\\([a-zA-Z]+)", r"\1", text)
    
    # Clean up double spaces
    text = text.replace("  ", " ")
    
    return text.strip()


def reserialize_question(q_text: str, new_num: int) -> str:
    """
    Replace original question number with new sequential number.
    Preserves sub-part markers like (i), (ii), (iii), (iv), (v).
    
    Args:
        q_text: Original question text starting with "1. " or "12. (ii)"
        new_num: New sequential number to assign
        
    Returns:
        Question text with updated number (preserving sub-parts)
    """
    # Pattern to match question number, optionally followed by sub-part marker
    # e.g., "12. (ii)" -> captures "(ii)" as sub_part
    pattern = r'^(SQ)?\d+\.\s*(\([ivx]+\)\s*)?'
    match = re.match(pattern, q_text)
    
    if match:
        sub_part = match.group(2) or ""  # e.g., "(ii) " or ""
        remaining_text = q_text[match.end():]
        return f"{new_num}. {sub_part}{remaining_text}"
    
    return q_text


def get_original_question_number(q_text: str) -> str:
    """
    Extract original question number from question text.
    
    Args:
        q_text: Question text starting with number
        
    Returns:
        Original question number (e.g., "1", "12. (ii)", "84. (iii)")
    """
    # Match question number with optional sub-part like "12. (ii)"
    match = re.match(r'^(SQ)?\d+\.?\s*(\([ivx]+\))?', q_text)
    return match.group(0).strip() if match else "?"


def parse_meta_info(meta: str) -> tuple[str, str]:
    """
    Parse meta string into year and marks components.
    
    Args:
        meta: Meta string in format "YEAR | MARKS"
        
    Returns:
        Tuple of (year_text, marks_text)
    """
    parts = meta.split('|')
    year = parts[0].strip() if len(parts) > 0 else ""
    marks = parts[1].strip() if len(parts) > 1 else ""
    return year, marks


# =============================================================================
# QUESTION DATA (imported from separate module)
# =============================================================================

from questions_data import QuestionData, questions_data


# =============================================================================
# SLIDE CREATION
# =============================================================================

class SlideBuilder:
    """Builder class for creating styled presentation slides."""
    
    def __init__(self, presentation: Presentation):
        self.prs = presentation
        self.slide_width = presentation.slide_width
        self.slide_height = presentation.slide_height
        self.content_height = int(self.slide_height * LAYOUT.content_ratio)
        self.margin = Inches(LAYOUT.margin)
        self.text_width = self.slide_width - (2 * self.margin)
    
    def create_slide(self, question: QuestionData) -> None:
        """
        Create a single styled slide for a question.
        
        Args:
            question: Question data dictionary with 'q' and 'meta' keys
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._set_background(slide)
        year, marks = parse_meta_info(question['meta'])
        self._add_year_label(slide, year)
        self._add_marks_label(slide, marks)
        self._add_accent_line(slide)
        self._add_question_text(slide, question['q'])
        # Extract original question number for the badge
        original_q_num = get_original_question_number(question['q'])
        self._add_question_badge(slide, original_q_num)
    
    def _set_background(self, slide) -> None:
        """Set slide background color."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS.background
    
    def _add_year_label(self, slide, year_text: str) -> None:
        """Add year label to top-left of slide."""
        box = slide.shapes.add_textbox(
            self.margin, 
            Inches(LAYOUT.top_margin), 
            Inches(3), 
            Inches(LAYOUT.header_height)
        )
        para = box.text_frame.paragraphs[0]
        para.text = f"{year_text}"
        para.font.name = FONTS.name
        para.font.size = Pt(FONTS.meta_size)
        para.font.color.rgb = COLORS.year
        para.font.bold = True
    
    def _add_marks_label(self, slide, marks_text: str) -> None:
        """Add marks label to top-right of slide."""
        box = slide.shapes.add_textbox(
            self.slide_width - self.margin - Inches(2),
            Inches(LAYOUT.top_margin),
            Inches(2),
            Inches(LAYOUT.header_height)
        )
        para = box.text_frame.paragraphs[0]
        para.text = f"{marks_text}"
        para.font.name = FONTS.name
        para.font.size = Pt(FONTS.meta_size)
        para.font.color.rgb = COLORS.marks
        para.font.bold = True
        para.alignment = PP_ALIGN.RIGHT
    
    def _add_accent_line(self, slide) -> None:
        """Add decorative accent line below header."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.margin,
            Inches(LAYOUT.accent_line_y),
            self.text_width,
            Inches(LAYOUT.accent_line_height)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS.accent_line
        line.line.fill.background()
    
    def _add_question_text(self, slide, question_text: str) -> None:
        """Add question text with original number preserved."""
        question_top = Inches(LAYOUT.question_top)
        question_height = self.content_height - question_top - Inches(0.3)
        
        box = slide.shapes.add_textbox(
            self.margin, 
            question_top, 
            self.text_width, 
            question_height
        )
        box.text_frame.word_wrap = True
        
        para = box.text_frame.paragraphs[0]
        cleaned_text = clean_chemistry_text(question_text)
        para.text = cleaned_text  # Keep original question numbers as-is
        para.font.name = FONTS.name
        para.font.size = Pt(FONTS.question_size)
        para.font.color.rgb = COLORS.text
    
    def _add_question_badge(self, slide, original_q_num: str) -> None:
        """Add small question number badge at bottom-left."""
        box = slide.shapes.add_textbox(
            self.margin,
            self.slide_height - Inches(LAYOUT.badge_bottom),
            Inches(1.5),
            Inches(0.4)
        )
        para = box.text_frame.paragraphs[0]
        para.text = f"Q{original_q_num}"
        para.font.name = FONTS.name
        para.font.size = Pt(FONTS.badge_size)
        para.font.color.rgb = COLORS.question_badge
        para.font.bold = True


# =============================================================================
# PRESENTATION GENERATOR
# =============================================================================

def create_presentation(
    questions: List[QuestionData],
    output_filename: str = "Chemistry_PYQ_Presentation.pptx"
) -> None:
    """
    Generate a PowerPoint presentation from question data.
    
    Args:
        questions: List of question dictionaries
        output_filename: Name of output PPTX file
    """
    prs = Presentation()
    builder = SlideBuilder(prs)
    
    for question in questions:
        builder.create_slide(question)
    
    prs.save(output_filename)
    print(f"✅ Presentation saved as {output_filename} with {len(questions)} slides.")


# =============================================================================
# MAIN ENTRY POINT
# =============================================================================

def main() -> None:
    """Main entry point for generating the chemistry presentation."""
    print("\n📚 Chemistry PYQ Presentation Generator")
    print("=" * 45)
    create_presentation(questions_data)
    print("=" * 45)


if __name__ == "__main__":
    main()